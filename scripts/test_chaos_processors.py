#!/usr/bin/env python3
"""Test script for Aily Chaos content processors."""

import asyncio
import tempfile
from pathlib import Path

from aily.chaos.config import ChaosConfig
from aily.chaos.processors import (
    ImageProcessor,
    PDFProcessor,
    PPTXProcessor,
    TextProcessor,
    VideoProcessor,
)


async def test_text_processor():
    """Test text processor."""
    print("Testing TextProcessor...")

    with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as f:
        f.write("# Test Document\n\nThis is a test.")
        temp_path = Path(f.name)

    try:
        config = ChaosConfig()
        processor = TextProcessor(config)

        result = await processor.process(temp_path)
        assert result is not None
        assert "Test Document" in result.text
        assert result.title == "Test Document"
        print("  ✓ TextProcessor works")
    finally:
        temp_path.unlink()


async def test_image_processor():
    """Test image processor (requires API key)."""
    print("Testing ImageProcessor...")

    # Create a simple test image
    from PIL import Image

    with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as f:
        temp_path = Path(f.name)

    try:
        # Create simple image
        img = Image.new("RGB", (100, 100), color="red")
        img.save(temp_path)

        config = ChaosConfig()
        processor = ImageProcessor(config)

        # Just test loading, API call requires key
        assert processor.can_process(temp_path)
        print("  ✓ ImageProcessor can_process works")
    finally:
        if temp_path.exists():
            temp_path.unlink()


async def test_pdf_processor():
    """Test PDF processor."""
    print("Testing PDFProcessor...")

    # Try to find a test PDF
    test_paths = [
        Path.home() / "aily_chaos" / "test.pdf",
        Path("test.pdf"),
    ]

    config = ChaosConfig()
    processor = PDFProcessor(config)

    found = False
    for path in test_paths:
        if path.exists():
            found = True
            result = await processor.process(path)
            if result:
                print(f"  ✓ PDFProcessor works: {result.title}")
            break

    if not found:
        print("  ⚠ No test PDF found, skipping")


async def test_pptx_processor():
    """Test PowerPoint processor."""
    print("Testing PPTXProcessor...")

    try:
        from pptx import Presentation
    except ImportError:
        print("  ⚠ python-pptx not installed, skipping")
        return

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        temp_path = Path(f.name)

    try:
        # Create simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        prs.save(temp_path)

        config = ChaosConfig()
        processor = PPTXProcessor(config)

        result = await processor.process(temp_path)
        assert result is not None
        assert "Test Presentation" in result.text
        print("  ✓ PPTXProcessor works")
    finally:
        if temp_path.exists():
            temp_path.unlink()


async def test_video_processor():
    """Test video processor detection."""
    print("Testing VideoProcessor...")

    config = ChaosConfig()
    processor = VideoProcessor(config)

    # Just test can_process
    test_path = Path("test.mp4")
    assert processor.can_process(test_path)
    print("  ✓ VideoProcessor can_process works")


async def main():
    """Run all tests."""
    print("=" * 50)
    print("Aily Chaos Processor Tests")
    print("=" * 50)

    await test_text_processor()
    await test_image_processor()
    await test_pdf_processor()
    await test_pptx_processor()
    await test_video_processor()

    print("=" * 50)
    print("Tests complete!")


if __name__ == "__main__":
    asyncio.run(main())
