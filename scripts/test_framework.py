#!/usr/bin/env python3
"""Unified ad hoc test framework for Aily.

This module consolidates the overlapping DIKIWI/Chaos/Kimi/manual test
scripts into a scenario-driven API that can be invoked from one CLI.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import random
import re
import shutil
import tempfile
import textwrap
import time
import traceback
from contextlib import contextmanager
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterator

from aily.browser.manager import BrowserUseManager
from aily.chaos.config import ChaosConfig
from aily.chaos.dikiwi_bridge import ChaosDikiwiBridge
from aily.chaos.processors import ImageProcessor, PDFProcessor, PPTXProcessor, TextProcessor, VideoProcessor
from aily.chaos.types import ExtractedContentMultimodal
from aily.config import SETTINGS
from aily.gating.drainage import RainDrop, RainType
from aily.graph.db import GraphDB
from aily.llm.client import LLMClient
from aily.llm.provider_routes import PrimaryLLMRoute
from aily.processing.router import ProcessingRouter
from aily.bot.message_intent import IntentRouter
from aily.push.feishu import FeishuPusher
from aily.queue.db import QueueDB
from aily.parser import registry
from aily.parser.parsers import parse_kimi
from aily.browser.fetcher import BrowserFetcher
from aily.processing.atomicizer import AtomicNoteGenerator
from aily.sessions.dikiwi_mind import DikiwiMind, DikiwiStage
from aily.sessions.entrepreneur_scheduler import EntrepreneurScheduler
from aily.sessions.reactor_scheduler import ReactorScheduler
from aily.sessions.reactor_scheduler import NozzleConfig
from aily.thinking.models import KnowledgePayload
from aily.thinking.orchestrator import ThinkingOrchestrator
from aily.ui.events import ui_event_hub
from aily.verify.evidence import EvidenceRun
from aily.writer.dikiwi_obsidian import DikiwiObsidianWriter

logger = logging.getLogger(__name__)

CHAOS_FOLDER = Path.home() / "aily_chaos"
DEFAULT_VAULT_PATH = Path(SETTINGS.dikiwi_vault_path or SETTINGS.obsidian_vault_path).expanduser()
DEFAULT_GRAPH_DB_PATH = DEFAULT_VAULT_PATH / ".aily" / "graph.db"
DEFAULT_LOG_DIR = Path(__file__).resolve().parent.parent / "logs" / "tests"
DEFAULT_RUN_DIR = Path(__file__).resolve().parent.parent / "logs" / "runs"

TEST_MESSAGES = [
    "【转向AI芯片架构的路径与优势 - Monica AI Chat】https://monica.im/share/chat?shareId=1jB54WO31xDzAIjL",
    "【模型的8bit和4bit量化原理与影响 - Monica AI Chat】https://monica.im/share/chat?shareId=VSvhr187W10wmQ5m",
    "【EDA软件的MCP蒸馏讨论 - Monica AI Chat】https://monica.im/share/chat?shareId=GGz9X6A7mnNeMqAJ",
    "【什么是MCP及其开发方法 - Monica AI Chat】https://monica.im/share/chat?shareId=9kH3k9l1jAPKh6t2",
    "【破除都灵裹尸布的迷雾 - Monica AI Chat】https://monica.im/share/chat?shareId=s36KOgwdvpjFZaEf",
    "【具有里程碑意义的AI技术 - Monica AI Chat】https://monica.im/share/chat?shareId=emlaeMyPoBUaFFfo",
    "【基于命令行工具转为MCP的可行性与借鉴工作 - Monica AI Chat】https://monica.im/share/chat?shareId=fdxLkrA92foijyIl",
    "【评估NVIDIA生成式AI技术用于EDA领域TCL脚本生成的适用性 - Monica AI Chat】https://monica.im/share/chat?shareId=BsA0KcdiGWQo4l09",
    "【PDK 评价体系与工艺线平衡分析 - Monica AI Chat】https://monica.im/share/chat?shareId=nLsKxwTCySW0p6Z3",
    "【芯片signoff规则制定方法论及学习资料 - Monica AI Chat】https://monica.im/share/chat?shareId=4cxQomLr6VD28Ofx",
]

ARMY_SAMPLE_CONTENT = """
CogniChip is a startup using AI to revolutionize chip design.
Founded by Faraj Aalaei, the company has raised $93M from NVIDIA and Intel.
Their ACI platform reduces chip design costs by 75% and time by 50%.
The team includes experts from Amazon, Google, Apple, and Synopsys.
They plan to tape out their first AI-designed chip by end of 2026.
The technology embeds physical constraints into AI models.
Their vision is to democratize chip design.
""".strip()


@dataclass
class RuntimeBundle:
    graph_db: GraphDB
    llm_client: LLMClient
    obsidian_writer: DikiwiObsidianWriter
    dikiwi_mind: DikiwiMind
    bridge: ChaosDikiwiBridge


def require_primary_api_key() -> str:
    route = PrimaryLLMRoute.resolve_route(SETTINGS, workload="default")
    if not route.api_key:
        raise RuntimeError("Set the provider API key for the configured default route before running this scenario.")
    return route.api_key


def build_primary_llm_client() -> LLMClient:
    return PrimaryLLMRoute.from_settings(SETTINGS, workload="default")


def clean_generated_vault(vault_path: Path, graph_db_path: Path | None = None) -> dict[str, Any]:
    removed_md = 0
    removed_dirs = 0
    generated_dirs = [
        vault_path / "00-Chaos",
        vault_path / "01-Data",
        vault_path / "02-Information",
        vault_path / "03-Knowledge",
        vault_path / "04-Insight",
        vault_path / "05-Wisdom",
        vault_path / "06-Impact",
        vault_path / "07-Proposal",
        vault_path / "08-Entrepreneurship",
        vault_path / "10-Knowledge",
        vault_path / "20-Innovation",
        vault_path / "30-Business",
        vault_path / "DIKIWI",
    ]
    for directory in generated_dirs:
        if directory.exists():
            for subpath in directory.rglob("*"):
                if subpath.is_file() and subpath.suffix == ".md":
                    subpath.unlink()
                    removed_md += 1
            for subpath in sorted(directory.rglob("*"), key=lambda p: len(p.parts), reverse=True):
                if subpath.is_dir() and not any(subpath.iterdir()):
                    subpath.rmdir()
                    removed_dirs += 1

    graph_reset = False
    graph_db_path = graph_db_path or (vault_path / ".aily" / "graph.db")
    if graph_db_path.exists():
        try:
            graph_db_path.unlink()
            graph_reset = True
        except PermissionError:
            logger.warning("Cannot unlink %s — may be locked by another process", graph_db_path)

    return {
        "removed_md": removed_md,
        "removed_dirs": removed_dirs,
        "graph_reset": graph_reset,
    }


def get_vault_status(vault_path: Path) -> dict[str, int]:
    if not vault_path.exists():
        return {}
    status: dict[str, int] = {}
    for subdir in sorted(vault_path.iterdir()):
        if subdir.is_dir() and not subdir.name.startswith("."):
            status[subdir.name] = len(list(subdir.rglob("*.md")))
    return status


def get_dir_samples(vault_path: Path, dir_name: str, limit: int = 3) -> list[tuple[str, int]]:
    d = vault_path / dir_name
    if not d.exists():
        return []
    files = sorted(d.rglob("*.md"), key=lambda p: p.stat().st_mtime, reverse=True)
    return [(str(f.relative_to(vault_path)), f.stat().st_size) for f in files[:limit]]


@contextmanager
def llm_trace_logging(log_path: Path | None) -> Iterator[Path | None]:
    if log_path is None:
        yield None
        return

    log_path.parent.mkdir(parents=True, exist_ok=True)
    original = LLMClient._chat_once

    async def patched(self, messages, temperature, response_format):
        call_record = {
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "model": self.model,
            "temperature": temperature,
            "response_format": response_format is not None,
            "messages": messages,
        }
        try:
            result = await original(self, messages, temperature, response_format)
            call_record["response"] = result
            call_record["success"] = True
        except Exception as exc:
            call_record["error"] = str(exc)
            call_record["success"] = False
            raise
        finally:
            with open(log_path, "a", encoding="utf-8") as f:
                f.write(json.dumps(call_record, ensure_ascii=False) + "\n")
        return result

    LLMClient._chat_once = patched
    try:
        yield log_path
    finally:
        LLMClient._chat_once = original


async def build_runtime(
    vault_path: Path = DEFAULT_VAULT_PATH,
    *,
    enable_business: bool = False,
) -> RuntimeBundle:
    graph_db = GraphDB(db_path=vault_path / ".aily" / "graph.db")
    await graph_db.initialize()

    llm_client = build_primary_llm_client()
    obsidian_writer = DikiwiObsidianWriter(vault_path=vault_path)
    dikiwi_mind = DikiwiMind(
        graph_db=graph_db,
        llm_client=llm_client,
        dikiwi_obsidian_writer=obsidian_writer,
    )

    if enable_business:
        nozzle_config = NozzleConfig(
            min_confidence=SETTINGS.minds.proposal_min_confidence,
            max_proposals_per_session=SETTINGS.minds.proposal_max_per_session,
        )
        dikiwi_mind.reactor_scheduler = ReactorScheduler(
            graph_db=graph_db,
            llm_client=llm_client,
            obsidian_writer=obsidian_writer,
            nozzle_config=nozzle_config,
            method_timeout_seconds=SETTINGS.reactor_method_timeout_seconds,
        )
        dikiwi_mind.entrepreneur_scheduler = EntrepreneurScheduler(
            graph_db=graph_db,
            llm_client=llm_client,
            obsidian_writer=obsidian_writer,
            proposal_min_confidence=SETTINGS.minds.proposal_min_confidence,
            proposal_max_per_session=SETTINGS.minds.proposal_max_per_session,
        )

    bridge = ChaosDikiwiBridge(
        dikiwi_mind=dikiwi_mind,
        processed_folder=CHAOS_FOLDER / ".processed",
    )
    return RuntimeBundle(
        graph_db=graph_db,
        llm_client=llm_client,
        obsidian_writer=obsidian_writer,
        dikiwi_mind=dikiwi_mind,
        bridge=bridge,
    )


async def close_runtime(runtime: RuntimeBundle) -> None:
    await runtime.graph_db.close()


async def scenario_processors() -> dict[str, Any]:
    results: dict[str, Any] = {}

    with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as f:
        f.write("# Test Document\n\nThis is a test.")
        temp_md = Path(f.name)

    try:
        config = ChaosConfig()
        text_processor = TextProcessor(config)
        text_result = await text_processor.process(temp_md)
        results["text"] = bool(text_result and "Test Document" in text_result.text)
    finally:
        temp_md.unlink(missing_ok=True)

    from PIL import Image

    with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as f:
        temp_jpg = Path(f.name)
    try:
        Image.new("RGB", (100, 100), color="red").save(temp_jpg)
        image_processor = ImageProcessor(ChaosConfig())
        results["image_can_process"] = image_processor.can_process(temp_jpg)
    finally:
        temp_jpg.unlink(missing_ok=True)

    try:
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_pptx = Path(f.name)
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Presentation"
            prs.save(temp_pptx)
            pptx_result = await PPTXProcessor(ChaosConfig()).process(temp_pptx)
            results["pptx"] = bool(pptx_result and "Test Presentation" in pptx_result.text)
        finally:
            temp_pptx.unlink(missing_ok=True)
    except ImportError:
        results["pptx"] = "skipped"

    results["video_can_process"] = VideoProcessor(ChaosConfig()).can_process(Path("test.mp4"))

    test_pdf_candidates = [CHAOS_FOLDER / "test.pdf", Path("test.pdf")]
    pdf_result = "skipped"
    for candidate in test_pdf_candidates:
        if candidate.exists():
            extracted = await PDFProcessor(ChaosConfig()).process(candidate)
            pdf_result = bool(extracted)
            break
    results["pdf"] = pdf_result
    return results


async def scenario_dikiwi_smoke(limit: int = 2, url: str | None = None) -> dict[str, Any]:
    require_primary_api_key()
    llm_resolver = PrimaryLLMRoute.build_settings_resolver(SETTINGS)
    client = llm_resolver("default")
    direct_checks = {
        "chat": False,
        "json": False,
        "classification": False,
    }

    try:
        response = await client.chat(
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": "What is the capital of France?"},
            ],
            temperature=0.5,
        )
        direct_checks["chat"] = "paris" in response.lower()
    except Exception:
        direct_checks["chat"] = False

    try:
        result = await client.chat_json(
            messages=[
                {"role": "system", "content": "Return JSON only."},
                {"role": "user", "content": "Return {\"ok\": true}"},
            ],
        )
        direct_checks["json"] = isinstance(result, dict)
    except Exception:
        direct_checks["json"] = False

    try:
        result = await client.chat_json(
            messages=[
                {"role": "system", "content": "Classify the topic and return JSON only."},
                {
                    "role": "user",
                    "content": (
                        'Classify "Reinforcement learning improves model alignment" into one of '
                        '["technology","business","science"] and return {"category":"..."}'
                    ),
                },
            ],
        )
        direct_checks["classification"] = isinstance(result, dict) and bool(result.get("category"))
    except Exception:
        direct_checks["classification"] = False

    db_dir = Path(tempfile.gettempdir()) / "aily_test"
    db_dir.mkdir(exist_ok=True)
    graph_db = GraphDB(db_path=db_dir / "test_dikiwi.db")
    await graph_db.initialize()

    try:
        mind = DikiwiMind(
            graph_db=graph_db,
            enabled=True,
            llm_client=llm_resolver("dikiwi"),
            llm_client_resolver=llm_resolver,
        )
        selected_messages = [url] if url else TEST_MESSAGES[:limit]
        pipeline_runs = []
        for index, msg in enumerate(selected_messages, 1):
            drop = RainDrop(
                rain_type=RainType.CHAT,
                id=f"smoke_{index}",
                content=msg,
                source="run_test_suite",
            )
            try:
                result = await mind.process_input(drop)
                pipeline_runs.append({
                    "ok": True,
                    "pipeline_id": result.pipeline_id,
                    "final_stage": result.final_stage_reached.name if result.final_stage_reached else None,
                    "stage_count": len(result.stage_results),
                    "total_time_ms": result.total_time_ms,
                })
            except Exception as exc:
                pipeline_runs.append({
                    "ok": False,
                    "error": str(exc),
                })
        metrics = mind.get_metrics()
    finally:
        await graph_db.close()

    return {
        "direct_checks": direct_checks,
        "pipeline_runs": pipeline_runs,
        "metrics": metrics,
    }


async def extract_pdf(pdf_path: Path) -> ExtractedContentMultimodal | None:
    return await PDFProcessor(ChaosConfig()).process(pdf_path)


async def extract_images(image_paths: list[Path]) -> list[ExtractedContentMultimodal]:
    processor = ImageProcessor(ChaosConfig())
    results = []
    for path in image_paths:
        extracted = await processor.process(path)
        if extracted:
            results.append(extracted)
    if not results:
        return []
    merged_text = "\n\n".join(r.text for r in results if r.text)
    merged = ExtractedContentMultimodal(
        text=merged_text,
        title=f"Image Session ({len(results)} photos)",
        source_type="image_session",
        source_path=image_paths[0],
        processing_method="multimodal-vision-session",
        metadata={"image_count": len(results), "paths": [str(p) for p in image_paths]},
        tags=[],
    )
    return [merged]


async def scenario_chaos_e2e(
    *,
    run_pdf: bool,
    run_images: bool,
    dry_run: bool,
    n_images: int,
    vault_path: Path = DEFAULT_VAULT_PATH,
) -> dict[str, Any]:
    items: list[tuple[str, ExtractedContentMultimodal]] = []

    if run_pdf:
        pdfs = sorted((CHAOS_FOLDER / "pdf").glob("*.pdf"))
        if pdfs:
            extracted = await extract_pdf(pdfs[0])
            if extracted:
                items.append((pdfs[0].name, extracted))

    if run_images:
        image_paths = sorted((CHAOS_FOLDER / "image").glob("*.jpg"))[:n_images]
        for extracted in await extract_images(image_paths):
            items.append((extracted.title or "image-session", extracted))

    if dry_run:
        return {
            "mode": "dry_run",
            "items": [
                {
                    "label": label,
                    "title": content.title,
                    "processing_method": content.processing_method,
                    "text_chars": len(content.text or ""),
                }
                for label, content in items
            ],
        }

    runtime = await build_runtime(vault_path=vault_path)
    try:
        results = []
        for label, content in items:
            bridge_result = await runtime.bridge.process_extracted_content(content)
            results.append({"label": label, **bridge_result})
        return {
            "items": results,
            "vault_status": get_vault_status(vault_path),
        }
    finally:
        await close_runtime(runtime)


async def scenario_url_audit(limit: int = 10, save_dir: Path | None = None) -> dict[str, Any]:
    save_dir = save_dir or Path.cwd()
    save_dir.mkdir(parents=True, exist_ok=True)

    classification_results = []
    for idx, msg in enumerate(TEST_MESSAGES[:limit], 1):
        intent = IntentRouter.analyze(msg)
        classification_results.append({
            "msg_num": idx,
            "intent_type": intent.intent_type.name,
            "url": intent.url,
            "confidence": intent.confidence,
            "reasoning": intent.reasoning,
        })

    browser = BrowserUseManager()
    await browser.start()
    fetch_results = []
    try:
        from aily.processing.markdownize import MarkdownizeProcessor

        processor = MarkdownizeProcessor(browser_manager=browser)
        for idx, msg in enumerate(TEST_MESSAGES[:limit], 1):
            url = IntentRouter.analyze(msg).url
            if not url:
                fetch_results.append({"msg_num": idx, "status": "no_url"})
                continue
            try:
                md_content = await processor.process_url(url, use_browser=True)
                output_file = save_dir / f"monica_msg_{idx}.md"
                output_file.write_text(md_content.markdown, encoding="utf-8")
                fetch_results.append({
                    "msg_num": idx,
                    "url": url,
                    "status": "success" if md_content.markdown else "error",
                    "source_type": md_content.source_type,
                    "title": md_content.title,
                    "text_length": len(md_content.markdown or ""),
                })
            except Exception as exc:
                fetch_results.append({
                    "msg_num": idx,
                    "url": url,
                    "status": "exception",
                    "error": str(exc),
                })
    finally:
        await browser.stop()

    useful_content = sum(
        1 for r in fetch_results if r.get("status") == "success" and r.get("text_length", 0) > 500
    )
    return {
        "classification_results": classification_results,
        "fetch_results": fetch_results,
        "summary": {
            "total_messages": limit,
            "fetch_successes": sum(1 for r in fetch_results if r.get("status") == "success"),
            "useful_content": useful_content,
        },
    }


async def scenario_full_pipeline(
    *,
    max_pdfs: int = 20,
    no_clean: bool = False,
    log_llm: bool = False,
    vault_path: Path = DEFAULT_VAULT_PATH,
    report_dir: Path | None = None,
    source_seed: int = 260502,
    phase_timeout_seconds: float = 600.0,
    force_business: bool = False,
) -> dict[str, Any]:
    report_dir = report_dir or DEFAULT_LOG_DIR / "e2e"
    report_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    report_file = report_dir / f"e2e_report_{timestamp}.json"
    llm_log = report_dir / f"llm_calls_{timestamp}.jsonl" if log_llm else None

    if not no_clean:
        clean_generated_vault(vault_path)

    pdf_dir = CHAOS_FOLDER / "pdf"
    pdf_files = sorted(pdf_dir.glob("*.pdf"))
    random.Random(source_seed).shuffle(pdf_files)
    pdf_files = pdf_files[:max_pdfs]
    if not pdf_files:
        return {"error": "no pdfs"}

    evidence = EvidenceRun(
        root_dir=DEFAULT_RUN_DIR,
        scenario=f"full_pipeline_{len(pdf_files)}pdf",
        vault_path=vault_path,
        graph_db_path=vault_path / ".aily" / "graph.db",
        source_paths=pdf_files,
        source_selector="seeded_random",
        source_seed=source_seed,
        mocked=False,
        fake_components=[],
        real_files=True,
        real_graph_db=True,
        real_vault=True,
        real_llm=True,
    )
    evidence.capture_before()
    progress: list[dict[str, Any]] = []

    def _record_progress(phase: str, state: str, **payload: Any) -> None:
        progress.append({
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "phase": phase,
            "state": state,
            **payload,
        })
        evidence.write_json("progress.json", progress)

    async def _run_phase(phase: str, awaitable: Any) -> Any:
        _record_progress(phase, "started", timeout_seconds=phase_timeout_seconds)
        started = time.monotonic()
        try:
            value = await asyncio.wait_for(awaitable, timeout=phase_timeout_seconds)
        except asyncio.TimeoutError as exc:
            elapsed = round(time.monotonic() - started, 2)
            _record_progress(phase, "timeout", elapsed_seconds=elapsed)
            raise TimeoutError(f"Phase {phase!r} exceeded {phase_timeout_seconds}s") from exc
        except Exception as exc:
            elapsed = round(time.monotonic() - started, 2)
            _record_progress(phase, "failed", elapsed_seconds=elapsed, error=str(exc))
            raise
        elapsed = round(time.monotonic() - started, 2)
        _record_progress(phase, "completed", elapsed_seconds=elapsed)
        return value

    try:
        with llm_trace_logging(llm_log):
            runtime = await build_runtime(vault_path=vault_path, enable_business=True)
            try:
                processor = PDFProcessor(config=ChaosConfig())
                doc_results = []
                total_start = time.monotonic()
                transcript_dir = vault_path / "00-Chaos"
                transcript_dir.mkdir(parents=True, exist_ok=True)
                semaphore = asyncio.Semaphore(max(1, SETTINGS.mineru_batch_extract_concurrency))

                async def _extract_one(pdf_path: Path) -> tuple[Path, ExtractedContentMultimodal | None, float]:
                    doc_start = time.monotonic()
                    async with semaphore:
                        extracted = await processor.process(pdf_path)
                    if extracted:
                        transcript_path = transcript_dir / f"{pdf_path.stem}.md"
                        transcript_path.write_text(
                            f"# {extracted.title or pdf_path.stem}\n\n{extracted.get_full_text()}",
                            encoding="utf-8",
                        )
                        # Copy MinerU images to vault so markdown references resolve
                        mineru_out = extracted.metadata.get("mineru_output_dir", "")
                        src_images = Path(mineru_out) / "images" if mineru_out else None
                        if src_images and src_images.is_dir():
                            dest_images = transcript_dir / "images"
                            dest_images.mkdir(parents=True, exist_ok=True)
                            for img in src_images.iterdir():
                                if img.is_file() and not (dest_images / img.name).exists():
                                    shutil.copy2(img, dest_images / img.name)
                    return pdf_path, extracted, round(time.monotonic() - doc_start, 2)

                extracted_results = await _run_phase(
                    "extract",
                    asyncio.gather(*[_extract_one(pdf_path) for pdf_path in pdf_files]),
                )
                extracted_contents: list[ExtractedContentMultimodal] = []
                extraction_elapsed: dict[str, float] = {}
                for pdf_path, extracted, elapsed in extracted_results:
                    extraction_elapsed[pdf_path.name] = elapsed
                    if extracted is None:
                        doc_results.append({"pdf": pdf_path.name, "error": "extraction_failed", "elapsed": elapsed})
                    else:
                        extracted_contents.append(extracted)

                bridge_batch = await _run_phase(
                    "dikiwi_batch",
                    runtime.bridge.process_extracted_content_batch(extracted_contents),
                )
                bridge_by_source = {
                    Path(item["source_path"]).name: item
                    for item in bridge_batch.get("results", [])
                    if item.get("source_path")
                }

                for pdf_path in pdf_files:
                    if any(item.get("pdf") == pdf_path.name for item in doc_results):
                        continue
                    doc_results.append({
                        "pdf": pdf_path.name,
                        "bridge_result": bridge_by_source.get(pdf_path.name, {"error": "missing_batch_result"}),
                        "elapsed": extraction_elapsed.get(pdf_path.name, 0.0),
                        "tokens": runtime.llm_client.get_usage_stats(),
                    })

                pre_business_vault_status = get_vault_status(vault_path)
                has_impact_outputs = pre_business_vault_status.get("06-Impact", 0) > 0
                business_should_run = force_business or has_impact_outputs

                reactor_elapsed = None
                proposals = []
                business_skipped_reason = ""
                if runtime.dikiwi_mind.reactor_scheduler is not None and business_should_run:
                    reactor_start = time.monotonic()
                    context = await _run_phase(
                        "reactor_context",
                        runtime.dikiwi_mind.reactor_scheduler._gather_context(),
                    )
                    proposals = await _run_phase(
                        "reactor_evaluate",
                        runtime.dikiwi_mind.reactor_scheduler.evaluate_context(
                            context,
                            persist=True,
                            output=True,
                        ),
                    )
                    reactor_elapsed = round(time.monotonic() - reactor_start, 2)
                elif not business_should_run:
                    business_skipped_reason = "no_impact_outputs"
                    _record_progress(
                        "reactor_evaluate",
                        "skipped",
                        reason=business_skipped_reason,
                        force_business=force_business,
                    )

                entrepreneur_elapsed = None
                if runtime.dikiwi_mind.entrepreneur_scheduler is not None and business_should_run:
                    entrepreneur_start = time.monotonic()
                    await _run_phase(
                        "entrepreneur",
                        runtime.dikiwi_mind.entrepreneur_scheduler._run_session_wrapper(),
                    )
                    entrepreneur_elapsed = round(time.monotonic() - entrepreneur_start, 2)

                final_vault_status = get_vault_status(vault_path)
                total_elapsed = round(time.monotonic() - total_start, 2)
                result = {
                    "documents": len(doc_results),
                    "results": doc_results,
                    "elapsed_seconds": total_elapsed,
                    "source_seed": source_seed,
                    "batch_incremental_ratio": bridge_batch.get("incremental_ratio"),
                    "batch_incremental_threshold": bridge_batch.get("incremental_threshold"),
                    "batch_higher_order_triggered": bridge_batch.get("higher_order_triggered"),
                    "reactor_elapsed_seconds": reactor_elapsed,
                    "entrepreneur_elapsed_seconds": entrepreneur_elapsed,
                    "reactor_proposals": len(proposals),
                    "business_skipped_reason": business_skipped_reason,
                    "vault_status": final_vault_status,
                    "pre_business_vault_status": pre_business_vault_status,
                    "vault_samples": {
                        dir_name: get_dir_samples(vault_path, dir_name)
                        for dir_name in ["00-Chaos", "05-Wisdom", "07-Proposal", "08-Entrepreneurship"]
                    },
                    "token_usage": runtime.llm_client.get_usage_stats(),
                    "llm_log_file": str(llm_log) if llm_log else None,
                }
                failures = [
                    {"pdf": item.get("pdf", ""), "error": str(item.get("error") or item.get("bridge_result", {}).get("error"))}
                    for item in doc_results
                    if item.get("error") or item.get("bridge_result", {}).get("error")
                ]
                evidence.finalize(
                    exit_code=0,
                    result=result,
                    failures=failures,
                    llm_log_file=str(llm_log) if llm_log else None,
                    ui_events=ui_event_hub.recent_events(limit=10000),
                    repo_root=Path(__file__).resolve().parent.parent,
                )
                result["evidence_dir"] = str(evidence.path)
                result["evidence_manifest"] = str(evidence.path / "manifest.json")
                report_file.write_text(json.dumps(result, indent=2, default=str), encoding="utf-8")
                result["report_file"] = str(report_file)
                return result
            finally:
                await close_runtime(runtime)
    except Exception as exc:
        stderr_text = traceback.format_exc()
        evidence.finalize(
            exit_code=1,
            result={"error": str(exc)},
            failures=[{"error": str(exc), "traceback": stderr_text}],
            llm_log_file=str(llm_log) if llm_log else None,
            ui_events=ui_event_hub.recent_events(limit=10000),
            stderr_text=stderr_text,
            repo_root=Path(__file__).resolve().parent.parent,
        )
        raise


async def scenario_legacy_atomicizer(
    *,
    url: str,
    open_id: str = "",
    clean_content: bool = False,
) -> dict[str, Any]:
    queue_db = QueueDB(SETTINGS.queue_db_path)
    graph_db = GraphDB(SETTINGS.graph_db_path)
    await queue_db.initialize()
    await graph_db.initialize()

    llm = LLMClient(
        base_url=SETTINGS.llm_base_url,
        api_key=SETTINGS.llm_api_key,
        model=SETTINGS.llm_model,
    )
    pusher = FeishuPusher(SETTINGS.feishu_app_id, SETTINGS.feishu_app_secret) if open_id else None
    fetcher = BrowserFetcher()

    def clean_kimi_content(text: str) -> str:
        ui_patterns = [r"新建会话", r"⌘", r"K\n", r"网站\n", r"文档\n", r"PPT\n", r"表格\n"]
        for pattern in ui_patterns:
            text = re.sub(pattern, "", text)
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        return "\n\n".join(lines)

    registry.register(r"^https://kimi\.moonshot\.cn/share/", parse_kimi)
    registry.register(r"^https://www\.kimi\.com/share/", parse_kimi)

    try:
        raw_text = await fetcher.fetch(url)
        content = clean_kimi_content(raw_text) if clean_content else raw_text
        parsed = registry.parse(url, content)
        raw_log_id = await queue_db.insert_raw_log(url, source="legacy_atomicizer") or "legacy_atomicizer"
        atomicizer = AtomicNoteGenerator(llm, graph_db)
        notes = await atomicizer.atomize(
            content=parsed.markdown[:15000],
            source_url=url,
            raw_log_id=raw_log_id,
        )
        all_connections = []
        for note in notes[:3]:
            all_connections.extend(await atomicizer.suggest_connections(note) or [])
        if pusher:
            await pusher.send_message(
                open_id,
                f"Legacy atomicizer test complete. Notes={len(notes)} Connections={len(all_connections)}",
            )
        return {
            "title": parsed.title,
            "markdown_chars": len(parsed.markdown),
            "notes": len(notes),
            "connections": len(all_connections),
        }
    finally:
        await fetcher.stop()
        await queue_db.close()
        await graph_db.close()


async def scenario_army() -> dict[str, Any]:
    payload = KnowledgePayload(
        content=ARMY_SAMPLE_CONTENT,
        source_url="https://www.kimi.com/share/example",
        source_title="CogniChip Deep Research",
    )
    llm = LLMClient(
        base_url=SETTINGS.llm_base_url,
        api_key=SETTINGS.llm_api_key,
        model=SETTINGS.llm_model,
    )
    orchestrator = ThinkingOrchestrator(llm_client=llm)
    result = await orchestrator.think(
        payload,
        options={"output_format": "obsidian", "max_insights": 5},
    )
    return {
        "confidence": result.confidence_score,
        "frameworks": [fi.framework_type.value for fi in result.framework_insights],
        "insights": len(result.synthesized_insights),
        "top_insights": [insight.title for insight in result.top_insights[:3]],
    }
