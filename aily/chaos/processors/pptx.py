"""PowerPoint processor."""

from __future__ import annotations

import asyncio
import base64
import io
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from aily.chaos.processors.base import ContentProcessor
from aily.chaos.types import ExtractedContentMultimodal, VisualElement

logger = logging.getLogger(__name__)


class PPTXProcessor(ContentProcessor):
    """Process PowerPoint presentations."""

    async def process(self, file_path: Path) -> ExtractedContentMultimodal | None:
        """Process PowerPoint file."""
        logger.info("Processing PowerPoint: %s", file_path.name)

        try:
            # Read presentation
            prs = await asyncio.to_thread(Presentation, str(file_path))

            slides_content = []
            visual_elements = []
            speaker_notes = []

            for i, slide in enumerate(prs.slides, 1):
                # Extract slide text
                slide_text = self._extract_slide_text(slide)
                slides_content.append(f"### Slide {i}\n\n{slide_text}")

                # Extract speaker notes
                if slide.has_notes_slide and self.config.pptx.extract_speaker_notes:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        speaker_notes.append(f"### Slide {i} Notes\n\n{notes}")

                # Count shapes for metadata
                shape_count = len(slide.shapes)

                # Create visual element placeholder
                if self.config.pptx.extract_slide_images:
                    element = VisualElement(
                        element_id=f"slide_{i}",
                        element_type="slide",
                        description=f"Slide {i}: {slide_text[:100]}..." if len(slide_text) > 100 else f"Slide {i}: {slide_text}",
                        source_page=i,
                    )
                    visual_elements.append(element)

            # Build full text
            text_parts = []

            # Title slide detection
            title = file_path.stem.replace("_", " ").replace("-", " ").title()
            if slides_content:
                first_slide_text = slides_content[0].lower()
                lines = slides_content[0].split("\n")
                for line in lines:
                    line = line.strip()
                    if line and not line.startswith("#") and len(line) < 100:
                        title = line
                        break

            text_parts.append(f"# {title}\n")
            text_parts.append(f"**Slides:** {len(prs.slides)}\n")

            # Main content
            text_parts.append("\n".join(slides_content))

            # Speaker notes
            if speaker_notes:
                text_parts.append("\n\n## Speaker Notes\n")
                text_parts.append("\n\n".join(speaker_notes))

            text = "\n\n".join(text_parts)

            return ExtractedContentMultimodal(
                text=text,
                title=title,
                source_type="presentation",
                source_path=file_path,
                visual_elements=visual_elements,
                processing_method="python-pptx",
                metadata={
                    "slide_count": len(prs.slides),
                    "has_notes": len(speaker_notes) > 0,
                },
            )

        except Exception as e:
            logger.exception("Failed to process PowerPoint: %s", e)
            return None

    def _extract_slide_text(self, slide) -> str:
        """Extract text from a slide."""
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        return "\n\n".join(text_parts)

    def can_process(self, file_path: Path) -> bool:
        """Check if file is a PowerPoint."""
        return file_path.suffix.lower() == ".pptx"
